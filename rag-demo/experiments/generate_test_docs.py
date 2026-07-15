# -*- coding: utf-8 -*-
"""生成多格式测试文档到 doc/ 目录"""

import os
import sys
from pathlib import Path

# Fix Windows console encoding
if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8")

BASE_DIR = Path(r"d:\AI\aiLearning\oneWeek\day8~10\RAGDemo\doc")
BASE_DIR.mkdir(parents=True, exist_ok=True)


# ═══════════════════════════════════════════════════════════
# 1. test.pdf — 3页, 中文+英文正文
# ═══════════════════════════════════════════════════════════
def generate_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from io import BytesIO
    from pypdf import PdfReader, PdfWriter
    import platform

    # 尝试注册中文字体
    cn_font = "Helvetica"
    system = platform.system()
    font_paths = []
    if system == "Windows":
        font_paths = [
            r"C:\Windows\Fonts\simhei.ttf",
            r"C:\Windows\Fonts\simsun.ttc",
            r"C:\Windows\Fonts\msyh.ttc",
        ]
    for fp in font_paths:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont("CN", fp))
                cn_font = "CN"
                break
            except Exception:
                continue

    pages = [
        [
            "=== Page 1: Introduction ===",
            "",
            "This is a multilingual test document for RAG document parsing.",
            "本测试文档用于验证 RAG 系统对多格式文档的解析能力。",
            "",
            "LangChain provides a unified interface for document loading,",
            "splitting, and retrieval. 它支持多种文件格式的解析。",
            "",
            "Key concepts:",
            "1. Document Loading — 文档加载",
            "2. Text Splitting — 文本分割",
            "3. Embedding & Retrieval — 嵌入与检索",
            "",
            "In natural language processing, tokenization is the process of",
            "splitting text into smaller units. 分词是自然语言处理中的基础步骤。",
        ],
        [
            "=== Page 2: Architecture ===",
            "",
            "The system uses a strategy pattern to route different file formats",
            "to appropriate document loaders.",
            "系统使用策略模式将不同格式路由到相应的文档加载器。",
            "",
            "Supported formats include PDF, Word, Excel, PowerPoint,",
            "Markdown, and plain text files.",
            "支持的格式包括 PDF、Word、Excel、PPT、Markdown 和纯文本。",
            "",
            "Vector databases are a key component of RAG systems.",
            "向量数据库是 RAG 系统的核心组件之一。",
            "They convert text into vector representations for efficient",
            "similarity search. 它们将文本转换为向量表示，支持相似度搜索。",
        ],
        [
            "=== Page 3: Summary ===",
            "",
            "Common embedding models include text-embedding-ada-002,",
            "bge-large-zh-v1.5, and other open-source alternatives.",
            "常用的 Embedding 模型包括 text-embedding-ada-002、bge-large-zh-v1.5 等。",
            "",
            "Testing ensures that all document formats can be correctly loaded",
            "with proper page/paragraph segmentation and no character corruption.",
            "测试确保所有格式的文档能被正确加载，分段合理，中文字符无乱码。",
            "",
            "Future enhancements: OCR for scanned documents,",
            "incremental loading, caching mechanisms, and parallel parsing.",
            "未来改进方向：OCR 扫描件支持、增量加载、缓存机制和并行解析。",
            "",
            "Thank you for using the Multi-Format Document Loader!",
            "感谢使用多格式文档加载器！",
        ],
    ]

    writer = PdfWriter()
    for pg_texts in pages:
        buf = BytesIO()
        c = rl_canvas.Canvas(buf, pagesize=A4)
        width, height = A4
        c.setFont(cn_font, 11)
        y = height - 40
        for line in pg_texts:
            if y < 40:
                c.showPage()
                y = height - 40
                c.setFont(cn_font, 11)
            c.drawString(50, y, line)
            y -= 18
        c.save()
        buf.seek(0)
        reader = PdfReader(buf)
        for page in reader.pages:
            writer.add_page(page)

    out_path = BASE_DIR / "test.pdf"
    with open(out_path, "wb") as f:
        writer.write(f)
    print(f"  ✅ test.pdf -> {out_path}")


# ═══════════════════════════════════════════════════════════
# 2. test.docx — 含标题 + 多段正文
# ═══════════════════════════════════════════════════════════
def generate_docx():
    import docx
    from docx.shared import Pt

    doc = docx.Document()
    doc.add_heading('多格式文档解析测试报告', level=0)
    doc.add_paragraph(
        '本文档用于验证 LangChain 多格式文档加载器的功能。'
    )

    doc.add_heading('一、项目背景', level=1)
    doc.add_paragraph(
        '随着大语言模型（LLM）的快速发展，检索增强生成（RAG）'
        '技术成为连接外部知识库与生成能力的重要桥梋。'
        '在实际应用中，企业知识库通常包含多种格式的文档。'
    )

    doc.add_heading('二、技术架构', level=1)
    doc.add_paragraph(
        '本系统采用策略模式设计，通过文件扩展名自动选择对应的文档加载器。'
        '不同格式对应不同的解析策略，但对外暴露统一的接口。'
    )

    doc.add_heading('三、支持的格式', level=1)
    items = [
        ('PDF', '使用 pypdf / pypdfium2 引擎进行解析，保留页面结构信息'),
        ('Word (.docx)', '使用 docx2txt 提取纯文本内容'),
        ('Excel (.xlsx / .csv)', '使用 unstructured 库解析表格数据'),
        ('PPT (.pptx)', '使用 unstructured 库提取投影片文本'),
        ('Markdown (.md)', '直接使用 TextLoader 读取'),
        ('纯文本 (.txt)', 'UTF-8 编码读取'),
    ]
    for fmt, desc in items:
        doc.add_paragraph(f'{fmt}：{desc}', style='List Bullet')

    doc.add_heading('四、测试数据', level=1)
    doc.add_paragraph(
        '以下测试数据涵盖中英文混合内容、结构化表格、'
        '代码示例等多种类型，用于全面评估各加载器的解析效果。'
    )

    doc.add_heading('五、预期结果', level=1)
    doc.add_paragraph(
        '所有格式的文档应能被正确加载，'
        '页码/段落划分合理，中文字符无乱码，'
        '元数据（文件名、类型、加载时间）完整记录。'
    )

    out_path = BASE_DIR / "test.docx"
    doc.save(str(out_path))
    print(f"  ✅ test.docx -> {out_path}")


# ═══════════════════════════════════════════════════════════
# 3. test.xlsx — 2个工作表，带表头 + 多行数据
# ═══════════════════════════════════════════════════════════
def generate_xlsx():
    import openpyxl

    wb = openpyxl.Workbook()

    ws1 = wb.active
    ws1.title = "产品数据"
    ws1.append(["序号", "产品名称", "类别", "价格(元)", "库存", "备注"])
    ws1.append([1, "语言模型实战指南", "技术书籍", 89.00, 156, "畅销书"])
    ws1.append([2, "大模型应用开发", "技术书籍", 75.50, 89, ""])
    ws1.append([3, "向量数据库入门", "技术书籍", 68.00, 234, "新书"])
    ws1.append([4, "RAG 系统设计", "技术书籍", 95.00, 0, "缺货"])
    ws1.append([5, "Embedding 原理", "技术书籍", 55.00, 312, ""])
    ws1.append([6, "Prompt Engineering", "英文原版", 120.00, 45, "进口"])
    ws1.append([7, "知识图谱构建", "技术书籍", 82.00, 178, ""])
    ws1.append([8, "多模态 AI 导论", "技术书籍", 78.50, 67, "新品"])

    ws2 = wb.create_sheet("性能指标")
    ws2.append(["指标名称", "Q1", "Q2", "Q3", "Q4", "年度目标"])
    ws2.append(["文档加载成功率", "98.5%", "99.1%", "99.3%", "99.5%", "99%"])
    ws2.append(["平均解析耗时(ms)", "320", "285", "260", "245", "<300"])
    ws2.append(["中文识别准确率", "96.2%", "97.0%", "97.8%", "98.1%", "97%"])
    ws2.append(["内存占用(MB)", "180", "165", "155", "148", "<200"])
    ws2.append(["并发处理能力", "50", "80", "120", "150", ">100"])

    out_path = BASE_DIR / "test.xlsx"
    wb.save(str(out_path))
    print(f"  ✅ test.xlsx -> {out_path}")


# ═══════════════════════════════════════════════════════════
# 4. test.csv — 带表头、普通数据、少量空单元格
# ═══════════════════════════════════════════════════════════
def generate_csv():
    lines = [
        "序号,姓名,部门,技能,熟练程度,备注",
        "1,张三,算法组,LangChain,精通,",
        "2,李四,工程组,Python,熟细,骨干成员",
        "3,王五,算法组,RAG架构,熟细,",
        "4,赵六,工程组,向量数据库,入门,",
        "5,孙七,算法组,Embedding,精通,论文作者",
        "6,周八,产品组,Prompt设计,熟细,",
        "7,,工程组,Docker,熟细,转岗中",
        "8,吴九,算法组,多模态,入门,实习生",
        "9,郑十,工程组,FastAPI,熟细,后端负责人",
    ]
    out_path = BASE_DIR / "test.csv"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"  ✅ test.csv -> {out_path}")


# ═══════════════════════════════════════════════════════════
# 5. test.pptx — 3页 PPT，含标题 + 正文 + 文本框
# ═══════════════════════════════════════════════════════════
def generate_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Slide 1: 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "多格式文档解析方案"
    slide.placeholders[1].text = "LangChain 统一加载器设计与实现\n2026年7月"

    # Slide 2: 架构概览
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "系统架构概览"

    bullets = [
        "策略模式：不同格式 = 不同加载策略",
        "统一接口：load() / load_batch()",
        "自动路由：根据扩展名选择 Loader",
        "元数据增强：文件名、类型、加载时间",
        "可扩展：新增格式只需在 LOADERS 字典中添加",
    ]
    tf = slide2.placeholders[1].text_frame
    tf.clear()
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.space_after = Pt(6)

    # 添加文本框
    tb = slide2.shapes.add_textbox(Inches(5), Inches(3.5), Inches(3), Inches(2))
    ttf = tb.text_frame
    p = ttf.add_paragraph()
    p.text = "核心优势："
    p.font.bold = True
    p2 = ttf.add_paragraph()
    p2.text = "✅ 格式全覆盖"
    p3 = ttf.add_paragraph()
    p3.text = "✅ 代码简洁"
    p4 = ttf.add_paragraph()
    p4.text = "✅ 易于维护"

    # Slide 3: 总结
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "总结与展望"

    bullets3 = [
        "已支持 PDF / DOCX / XLSX / CSV / PPTX / MD / TXT 七种格式",
        "中文内容解析正常，无乱码问题",
        "下一步：支持 OCR 扫描件、HTML 网页、JSON 数据",
        "持续优化：增量加载、缓存机制、并行解析",
    ]
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    for i, txt in enumerate(bullets3):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = txt
        p.space_after = Pt(6)

    # 添加文本框
    tb2 = slide3.shapes.add_textbox(Inches(5), Inches(3.5), Inches(3), Inches(2))
    ttf2 = tb2.text_frame
    p = ttf2.add_paragraph()
    p.text = "感谢阅读！"
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = 1  # CENTER
    p2 = ttf2.add_paragraph()
    p2.text = "RAG Demo Project"
    p2.font.size = Pt(14)
    p2.alignment = 1

    out_path = BASE_DIR / "test.pptx"
    prs.save(str(out_path))
    print(f"  ✅ test.pptx -> {out_path}")


# ═══════════════════════════════════════════════════════════
# 6. test.md — 含标题、列表、代码块
# ═══════════════════════════════════════════════════════════
def generate_md():
    content = """# 多格式文档加载器 - 使用说明

## 概述

本模块实现了基于 LangChain 的统一文档加载器，支持多种常见文件格式的解析。

## 快速开始

```python
from multi_format_loader import MultiFormatDocumentLoader

loader = MultiFormatDocumentLoader(encoding="utf-8")

# 加载单个文件
docs = loader.load("document.pdf")

# 批量加载整个文件夹
all_docs = loader.load_batch("./documents/")
```

## 支持的格式

| 扩展名 | 加载器 | 说明 |
|--------|--------|------|
| .pdf | PyPDFLoader | 基于 pypdf 引擎 |
| .docx | Docx2txtLoader | 提取 Word 文本 |
| .xlsx | UnstructuredExcelLoader | 解析 Excel 表格 |
| .csv | UnstructuredExcelLoader | CSV 作为电子表格解析 |
| .pptx | UnstructuredPowerPointLoader | 提取投影片内容 |
| .md | TextLoader | Markdown 纯文本 |
| .txt | TextLoader | UTF-8 编码文本 |

## 代码示例

以下是加载器核心代码的结构：

```python
class MultiFormatDocumentLoader:
    LOADERS = {
        '.pdf': PyPDFLoader,
        '.docx': Docx2txtLoader,
        '.xlsx': UnstructuredExcelLoader,
        '.pptx': UnstructuredPowerPointLoader,
        '.md': TextLoader,
        '.txt': TextLoader,
        '.csv': UnstructuredExcelLoader,
    }

    def load(self, file_path: str) -> list:
        ext = Path(file_path).suffix.lower()
        loader_class = self.LOADERS.get(ext)
        loader = loader_class(file_path)
        return loader.load()
```

## 注意事项

1. 确保安装了所有依赖包
2. PDF 文件建议为文本型（非扫描图片）
3. Excel 文件会自动解析第一个工作表
4. 中文内容请确保文件编码为 UTF-8

## 参考链接

- [LangChain 文档加载器](https://python.langchain.com/docs/modules/data_connection/document_loaders/)
- [Unstructured 库文档](https://docs.unstructured.io/)
"""
    out_path = BASE_DIR / "test.md"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"  ✅ test.md -> {out_path}")


# ═══════════════════════════════════════════════════════════
# 7. test.txt — 中文 + 英文混合 UTF-8
# ═══════════════════════════════════════════════════════════
def generate_txt():
    content = """多格式文档解析测试 - 纯文本文件

本节介绍 LangChain 多格式文档加载器的基本功能。

LangChain is a framework for developing applications powered by large language models (LLMs).
It simplifies the process of building AI applications by providing reusable components.
朗格链是一个用于开发大型语言模型应用的框架。

文档加载器支持以下功能：
1. 自动检测文件格式并选择对应的加载器
2. 统一的输出格式（Document 对象列表）
3. 丰富的元数据支持（文件名、类型、加载时间）
4. 批量加载能力（一次加载整个文件夹）

In natural language processing, document parsing is a fundamental step.
文档解析是自然语言处理中的基础步骤。

The system uses a strategy pattern to route different file formats to appropriate loaders.
系统使用策略模式将不同格式路由到相应的加载器。

测试数据段落二：
向量数据库（Vector Database）是 RAG 系统的核心组件之一。
Vector databases are a key component of RAG systems.
它们将文本转换为向量表示，支持高效的相似度搜索。
They convert text into vector representations for efficient similarity search.

测试数据段落三：
Embedding 模型将非结构化文本转换为密集向量。
Embedding models convert unstructured text into dense vectors.
常用的 Embedding 模型包括 text-embedding-ada-002、bge-large-zh-v1.5 等。
Common embedding models include text-embedding-ada-002, bge-large-zh-v1.5, and others.

Thank you for using the Multi-Format Document Loader!
感谢使用多格式文档加载器！
"""
    out_path = BASE_DIR / "test.txt"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"  ✅ test.txt -> {out_path}")


# ═══════════════════════════════════════════════════════════
# 执行所有生成
# ═══════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("生成测试文档到:", BASE_DIR)
    print()
    generate_pdf()
    generate_docx()
    generate_xlsx()
    generate_csv()
    generate_pptx()
    generate_md()
    generate_txt()
    print()
    print("全部完成！")
